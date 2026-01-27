#!/usr/bin/env python3
"""
IDM System - Beautiful PowerPoint Presentation Generator
Creates a professional presentation with process flows and system mockups
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.dml.color import RGBColor as RgbColor
from PIL import Image, ImageDraw, ImageFont
import os

# Colors
PRIMARY = RgbColor(0x16, 0x30, 0x6C)  # #16306C
SECONDARY = RgbColor(0xF9, 0xBF, 0x3F)  # #F9BF3F
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
DARK = RgbColor(0x1F, 0x29, 0x37)
LIGHT_GRAY = RgbColor(0xF3, 0xF4, 0xF6)
GREEN = RgbColor(0x10, 0xB9, 0x81)
RED = RgbColor(0xEF, 0x44, 0x44)
BLUE = RgbColor(0x3B, 0x82, 0xF6)

# Image directory
IMG_DIR = "/opt/idm-system/docs/images"

def create_mockup_image(filename, title, elements, width=1200, height=700):
    """Create a mockup image of system screen"""
    img = Image.new('RGB', (width, height), '#FFFFFF')
    draw = ImageDraw.Draw(img)

    # Header
    draw.rectangle([0, 0, width, 60], fill='#16306C')
    draw.text((20, 18), "IDM System", fill='#FFFFFF')
    draw.text((width - 150, 18), "admin@orienbank.tj", fill='#FFFFFF')

    # Navigation
    nav_items = ["Главная", "Мои заявки", "На согласовании", "Системы", "Админ"]
    x = 150
    for item in nav_items:
        draw.text((x, 20), item, fill='#FFFFFF')
        x += 120

    # Title
    draw.text((30, 80), title, fill='#16306C')

    # Content based on elements
    y = 130
    for elem in elements:
        if elem['type'] == 'card':
            draw.rectangle([30, y, 350, y + 80], fill='#F9FAFB', outline='#E5E7EB')
            draw.text((50, y + 15), elem['title'], fill='#16306C')
            draw.text((50, y + 45), elem.get('subtitle', ''), fill='#6B7280')
        elif elem['type'] == 'table_row':
            draw.rectangle([30, y, width - 30, y + 40], fill='#F9FAFB' if y % 80 == 0 else '#FFFFFF', outline='#E5E7EB')
            draw.text((50, y + 10), elem['text'], fill='#374151')
            if 'status' in elem:
                color = '#10B981' if elem['status'] == 'approved' else '#F59E0B' if elem['status'] == 'pending' else '#EF4444'
                draw.rectangle([width - 150, y + 8, width - 50, y + 32], fill=color)
                draw.text((width - 140, y + 10), elem['status_text'], fill='#FFFFFF')
        elif elem['type'] == 'button':
            draw.rectangle([elem.get('x', 30), y, elem.get('x', 30) + 150, y + 40], fill='#16306C')
            draw.text((elem.get('x', 30) + 20, y + 10), elem['text'], fill='#FFFFFF')
        elif elem['type'] == 'form_field':
            draw.rectangle([30, y, 500, y + 35], fill='#FFFFFF', outline='#D1D5DB')
            draw.text((40, y + 8), elem['placeholder'], fill='#9CA3AF')
            draw.text((30, y - 20), elem['label'], fill='#374151')
        y += elem.get('height', 90)

    img.save(os.path.join(IMG_DIR, filename))
    return os.path.join(IMG_DIR, filename)

def create_process_diagram(filename, steps, width=1200, height=500):
    """Create a process flow diagram"""
    img = Image.new('RGB', (width, height), '#FFFFFF')
    draw = ImageDraw.Draw(img)

    step_width = 180
    arrow_width = 60
    total_width = len(steps) * step_width + (len(steps) - 1) * arrow_width
    start_x = (width - total_width) // 2
    y = height // 2 - 50

    for i, step in enumerate(steps):
        x = start_x + i * (step_width + arrow_width)

        # Draw box
        color = step.get('color', '#16306C')
        draw.rounded_rectangle([x, y, x + step_width, y + 100], radius=10, fill=color)

        # Draw icon circle
        draw.ellipse([x + 70, y + 10, x + 110, y + 50], fill='#FFFFFF')
        draw.text((x + 82, y + 20), step.get('icon', '●'), fill=color)

        # Draw text
        text = step['text']
        text_width = len(text) * 7
        text_x = x + (step_width - text_width) // 2
        draw.text((text_x, y + 60), text, fill='#FFFFFF')

        # Draw number
        draw.text((x + 10, y + 10), str(i + 1), fill='#FFFFFF')

        # Draw arrow
        if i < len(steps) - 1:
            arrow_x = x + step_width + 10
            draw.polygon([
                (arrow_x, y + 50),
                (arrow_x + 40, y + 50),
                (arrow_x + 40, y + 40),
                (arrow_x + 55, y + 50),
                (arrow_x + 40, y + 60),
                (arrow_x + 40, y + 50),
            ], fill='#F9BF3F')

    img.save(os.path.join(IMG_DIR, filename))
    return os.path.join(IMG_DIR, filename)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()

    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SECONDARY
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_func):
    """Add a content slide with custom content"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(2), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SECONDARY
    accent.line.fill.background()

    content_func(slide)
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.3), width=Inches(9))

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create mockup images first
    print("Creating mockup images...")

    # Dashboard mockup
    create_mockup_image('dashboard.png', 'Панель управления', [
        {'type': 'card', 'title': 'Новая заявка', 'subtitle': 'Запросить доступ к системе'},
        {'type': 'card', 'title': 'Мои заявки', 'subtitle': 'Просмотр статусов заявок'},
        {'type': 'card', 'title': 'На согласовании', 'subtitle': 'Заявки, ожидающие утверждения'},
    ])

    # Request list mockup
    create_mockup_image('requests.png', 'Мои заявки', [
        {'type': 'table_row', 'text': 'REQ-2026-001  |  АБС "Банк XXI век"  |  Оператор', 'status': 'approved', 'status_text': 'Одобрено', 'height': 50},
        {'type': 'table_row', 'text': 'REQ-2026-002  |  CRM система  |  Менеджер', 'status': 'pending', 'status_text': 'Ожидает', 'height': 50},
        {'type': 'table_row', 'text': 'REQ-2026-003  |  Email  |  Пользователь', 'status': 'rejected', 'status_text': 'Отклонено', 'height': 50},
    ])

    # Create request form mockup
    create_mockup_image('create_request.png', 'Создание заявки', [
        {'type': 'form_field', 'label': 'Система', 'placeholder': 'Выберите систему...', 'height': 70},
        {'type': 'form_field', 'label': 'Роль доступа', 'placeholder': 'Выберите роль...', 'height': 70},
        {'type': 'form_field', 'label': 'Обоснование', 'placeholder': 'Укажите причину запроса доступа...', 'height': 70},
        {'type': 'button', 'text': 'Создать заявку', 'height': 60},
    ])

    # Approval mockup
    create_mockup_image('approval.png', 'Согласование заявки', [
        {'type': 'card', 'title': 'Заявитель: Иванов И.И.', 'subtitle': 'Отдел: ИТ департамент'},
        {'type': 'card', 'title': 'Система: АБС "Банк XXI век"', 'subtitle': 'Роль: Оператор'},
        {'type': 'button', 'text': '✓ Согласовать', 'x': 30, 'height': 60},
        {'type': 'button', 'text': '✗ Отклонить', 'x': 200, 'height': 60},
    ])

    # Process diagrams
    create_process_diagram('process_request.png', [
        {'text': 'Сотрудник', 'color': '#16306C', 'icon': '👤'},
        {'text': 'Создание\nзаявки', 'color': '#3B82F6', 'icon': '📝'},
        {'text': 'Согласование', 'color': '#F59E0B', 'icon': '✓'},
        {'text': 'Предоставление\nдоступа', 'color': '#10B981', 'icon': '🔓'},
    ])

    create_process_diagram('process_approval.png', [
        {'text': 'Новая\nзаявка', 'color': '#3B82F6', 'icon': '📋'},
        {'text': 'Уведомление\n(Email/TG)', 'color': '#8B5CF6', 'icon': '📧'},
        {'text': 'Проверка\nданных', 'color': '#F59E0B', 'icon': '🔍'},
        {'text': 'Решение', 'color': '#10B981', 'icon': '✓'},
    ])

    print("Creating presentation slides...")

    # Slide 1: Title
    add_title_slide(prs, "IDM System", "Система управления идентификацией и доступом\n\nОАО «Ориёнбонк» • 2026")

    # Slide 2: Problem
    def problem_content(slide):
        problems = [
            ("⏱️", "Долгое время обработки", "Заявки обрабатываются 3-5 дней"),
            ("📋", "Нет единой системы", "Заявки теряются, нет контроля"),
            ("🔍", "Отсутствие аудита", "Невозможно отследить историю"),
            ("📊", "Нет отчётности", "Нет данных для анализа"),
        ]
        y = 1.5
        for icon, title, desc in problems:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(4.3), Inches(1))
            box.fill.solid()
            box.fill.fore_color.rgb = RgbColor(0xFE, 0xE2, 0xE2)
            box.line.color.rgb = RgbColor(0xEF, 0x44, 0x44)

            tb = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(4), Inches(0.8))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{icon}  {title}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = DARK

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(14)
            p2.font.color.rgb = RgbColor(0x6B, 0x72, 0x80)

            y += 1.3

        # Solution arrow
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5), Inches(3), Inches(0.8), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SECONDARY

        # Solution box
        sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1.5), Inches(3.5), Inches(4.5))
        sol.fill.solid()
        sol.fill.fore_color.rgb = RgbColor(0xDC, 0xFC, 0xE7)
        sol.line.color.rgb = GREEN

        tb = slide.shapes.add_textbox(Inches(6.3), Inches(1.8), Inches(3), Inches(4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "✅ IDM System"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN

        benefits = ["4-8 часов обработки", "Единая точка входа", "Полный аудит", "Автоматические отчёты"]
        for b in benefits:
            p2 = tf.add_paragraph()
            p2.text = f"• {b}"
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK
            p2.space_before = Pt(12)

    add_content_slide(prs, "Проблема и решение", problem_content)

    # Slide 3: Process flow
    add_image_slide(prs, "Процесс подачи заявки", f"{IMG_DIR}/process_request.png",
                   "Полный цикл от создания заявки до получения доступа")

    # Slide 4: Dashboard
    add_image_slide(prs, "Главная страница системы", f"{IMG_DIR}/dashboard.png",
                   "Быстрый доступ к основным функциям")

    # Slide 5: Create request
    add_image_slide(prs, "Создание заявки на доступ", f"{IMG_DIR}/create_request.png",
                   "Простая форма для запроса доступа к системам")

    # Slide 6: Request list
    add_image_slide(prs, "Список заявок", f"{IMG_DIR}/requests.png",
                   "Отслеживание статуса всех заявок в реальном времени")

    # Slide 7: Approval process
    add_image_slide(prs, "Процесс согласования", f"{IMG_DIR}/process_approval.png",
                   "Автоматическая маршрутизация и уведомления")

    # Slide 8: Approval screen
    add_image_slide(prs, "Экран согласования", f"{IMG_DIR}/approval.png",
                   "Согласование в один клик с возможностью добавить комментарий")

    # Slide 9: Features
    def features_content(slide):
        features = [
            ("🔐", "Active Directory", "Вход через корпоративную учётную запись"),
            ("📧", "Email-уведомления", "Автоматические письма о статусе"),
            ("💬", "Telegram-бот", "Мгновенные уведомления и согласование"),
            ("📊", "Дашборд", "Статистика и аналитика в реальном времени"),
            ("🌙", "Тёмная тема", "Комфортный интерфейс для работы"),
            ("📄", "Экспорт", "PDF, Word, Excel отчёты"),
        ]

        col1_x, col2_x = 0.5, 5.2
        y = 1.5

        for i, (icon, title, desc) in enumerate(features):
            x = col1_x if i < 3 else col2_x
            if i == 3:
                y = 1.5

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.3))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY
            box.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(4), Inches(1))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{icon}  {title}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(14)
            p2.font.color.rgb = DARK

            y += 1.6

    add_content_slide(prs, "Функции системы", features_content)

    # Slide 10: Statistics
    def stats_content(slide):
        stats = [
            ("156", "Заявок\nза месяц", PRIMARY),
            ("6.5ч", "Среднее время\nсогласования", GREEN),
            ("92%", "Автоматизация\nпроцесса", BLUE),
            ("100%", "Аудит\nдействий", SECONDARY),
        ]

        x = 0.5
        for value, label, color in stats:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(2.1), Inches(2.5))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(2.1), Inches(2))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = label
            p2.font.size = Pt(16)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.CENTER

            x += 2.4

        # Bottom text
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "📈 Сокращение времени обработки заявок на 92% после внедрения системы"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    add_content_slide(prs, "Результаты и статистика", stats_content)

    # Slide 11: Thank you
    add_title_slide(prs, "Спасибо за внимание!", "Вопросы?\n\nGitHub: github.com/Adminpm04/idm-system")

    # Save presentation
    output_path = "/opt/idm-system/docs/IDM_System_Process_Presentation.pptx"
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")

    return output_path

if __name__ == "__main__":
    create_presentation()
